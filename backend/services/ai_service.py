"""AI service supporting OpenAI and AWS Bedrock (Claude)."""
import json
import asyncio
from typing import List, Optional
from config import settings


def _get_openai_client():
    from openai import AsyncOpenAI
    return AsyncOpenAI(api_key=settings.OPENAI_API_KEY)


def _get_bedrock_client():
    import boto3
    return boto3.client(
        "bedrock-runtime",
        region_name=settings.AWS_REGION,
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
    )


async def _llm_call(prompt: str, system: str = "", model: str = "openai") -> str:
    """Unified LLM call supporting OpenAI and Bedrock."""
    provider = settings.LLM_PROVIDER if model == "openai" else model
    if not settings.OPENAI_API_KEY and "bedrock" not in provider:
        return "[AI features require OPENAI_API_KEY in .env]"
    try:
        if "bedrock" in provider or provider == "bedrock":
            # Use cross-region inference profile IDs (required for Claude 4+)
            if "opus" in provider:
                model_id = "apac.anthropic.claude-opus-4-20250514-v1:0"
            elif "sonnet" in provider or provider == "bedrock":
                model_id = settings.BEDROCK_MODEL_ID  # from .env
            else:
                model_id = settings.BEDROCK_MODEL_ID
            client = _get_bedrock_client()
            body = json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": settings.CHAT_MAX_TOKENS,
                "system": system,
                "messages": [{"role": "user", "content": prompt}],
            })
            response = client.invoke_model(modelId=model_id, body=body)
            result = json.loads(response["body"].read())
            return result["content"][0]["text"]
        else:
            client = _get_openai_client()
            messages = []
            if system:
                messages.append({"role": "system", "content": system})
            messages.append({"role": "user", "content": prompt})
            resp = await client.chat.completions.create(
                model=settings.OPENAI_MODEL,
                messages=messages,
                max_tokens=settings.CHAT_MAX_TOKENS,
                temperature=0.3,
            )
            return resp.choices[0].message.content
    except Exception as e:
        return f"[LLM error: {e}]"


async def generate_article_metadata(text: str, filename: str = "") -> dict:
    """Generate title, summaries, insights and tags from PDF text."""
    # Check we have a usable LLM provider
    provider = settings.LLM_PROVIDER
    has_provider = ("bedrock" in provider) or bool(settings.OPENAI_API_KEY)
    if not has_provider:
        return {
            "title": filename.replace(".pdf", "").replace("_", " "),
            "summary_en": "Summary not available - please add API key to .env",
            "summary_th": "ยังไม่มีสรุป — กรุณาเพิ่ม API key ในไฟล์ .env",
            "key_insights": ["Add OPENAI_API_KEY or configure Bedrock in .env"],
            "tags": ["steel", "manufacturing"],
        }
    system = (
        "You are an expert technical editor for a structural steel manufacturing company (SYS) in Thailand. "
        "You read technical articles and extract structured metadata. "
        "Respond ONLY with valid JSON."
    )
    prompt = f"""Analyze this technical article text and return a JSON object with these fields:
- title: string (clean article title, infer from content if needed)
- summary_en: string (2-3 paragraph English summary)
- summary_th: string (2-3 paragraph Thai summary)
- key_insights: array of 5-8 bullet point strings in English
- tags: array of 8-15 keyword strings (lowercase, relevant to steel industry)

Article filename: {filename}

Article text (first 8000 chars):
{text[:8000]}

Return ONLY the JSON object, no other text."""

    result = await _llm_call(prompt, system)
    try:
        # Strip markdown code blocks if present
        clean = result.strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        return json.loads(clean)
    except Exception:
        return {
            "title": filename.replace(".pdf", "").replace("_", " "),
            "summary_en": "Summary not available - please add API key to .env",
            "summary_th": "ยังไม่มีสรุป — กรุณาเพิ่ม API key ในไฟล์ .env",
            "key_insights": ["Add OPENAI_API_KEY to .env to generate insights"],
            "tags": ["steel", "manufacturing"],
        }


async def answer_question(question: str, context_docs: list, history: list = [],
                          model: str = "openai", web_results: list = []) -> dict:
    """RAG-based question answering with optional web search results."""
    context = "\n\n---\n\n".join([
        f"Article: {d.get('title', 'Unknown')}\n{d.get('content', '')}"
        for d in context_docs
    ])

    web_context = ""
    if web_results:
        web_context = "\n\n=== ONLINE SEARCH RESULTS (retrieved right now) ===\n" + "\n\n".join([
            f"Source [{i+1}]: {r.get('title','')}\nURL: {r.get('source_url','')}\nContent: {r.get('content','')}"
            for i, r in enumerate(web_results)
        ]) + "\n=== END OF SEARCH RESULTS ==="

    has_web = bool(web_results)
    system = (
        "You are a knowledgeable assistant for SYS (Structural Steel Manufacturing), Thailand. "
        + ("You have been provided with REAL-TIME online search results above. "
           "You MUST use these search results to answer the question. "
           "Cite the source titles and URLs in your answer. "
           if has_web else
           "Answer questions using the provided article context. ")
        + "Be concise and accurate. "
        "Respond in the same language as the question (Thai or English)."
    )
    history_text = "\n".join([f"{h['role'].upper()}: {h['content']}" for h in history[-6:]])

    prompt = f"""{"Internal Knowledge Base Articles:" if context else ""}
{context}
{web_context}

{"Previous conversation:" + chr(10) + history_text if history_text else ""}

Question: {question}

{"Use the online search results above to answer. Cite sources." if has_web else "Answer based on the articles above."}"""

    answer = await _llm_call(prompt, system, model)
    sources = [d.get("title", "") for d in context_docs[:3] if d.get("title")]
    web_sources = [{"title": r.get("title", ""), "url": r.get("source_url", "")} for r in web_results[:5]]
    return {"answer": answer, "sources": sources, "web_sources": web_sources}


async def generate_report_markdown(articles: list, title: str, language: str = "en") -> str:
    """Generate a comprehensive research report by generating each section separately."""
    n = len(articles)

    if language == "th":
        system = (
            "คุณคือนักเขียนรายงานเทคนิคมืออาชีพสำหรับบริษัท SYS (ผู้ผลิตเหล็กโครงสร้าง) ประเทศไทย "
            "เขียนรายงานเป็นภาษาไทยทั้งหมด ห้ามใช้ภาษาอังกฤษในเนื้อหา ยกเว้นคำศัพท์เทคนิคเฉพาะทาง "
            "เขียนให้ละเอียดและครบถ้วน"
        )
    else:
        system = (
            "You are a senior technical report writer for SYS steel manufacturing (Thailand). "
            "Write comprehensive, detailed, professional reports in English. "
            "Each section must be thorough with multiple paragraphs and specific technical details."
        )

    articles_text = "\n\n".join([
        f"### Article {i+1}: {a.title}\n"
        f"**Category:** {a.topic_category} | **Date:** {a.publication_date}\n"
        f"**Summary:** {a.summary_en if language == 'en' else (a.summary_th or a.summary_en)}\n"
        f"**Tags:** {', '.join(a.tags or [])}\n"
        f"**Key Insights:**\n" + "\n".join([f"- {ins}" for ins in (a.key_insights or [])])
        for i, a in enumerate(articles)
    ])

    # Generate each section separately for longer, more detailed output
    async def gen_section(section_prompt: str) -> str:
        return await _llm_call_report(section_prompt, system)

    if language == "th":
        sections = await asyncio.gather(
            gen_section(f"""เขียน "บทสรุปผู้บริหาร" สำหรับรายงานชื่อ "{title}" จากบทความเหล่านี้:
{articles_text}
เขียนอย่างน้อย 4 ย่อหน้า ครอบคลุมความสำคัญ ประเด็นหลัก และผลกระทบต่อ SYS"""),
            gen_section(f"""เขียน "ผลการวิจัยหลัก" สำหรับรายงานชื่อ "{title}" จากบทความเหล่านี้:
{articles_text}
วิเคราะห์แต่ละบทความอย่างละเอียด (2-3 ย่อหน้าต่อบทความ) แล้วสังเคราะห์ประเด็นร่วม"""),
            gen_section(f"""เขียน "แนวโน้มเทคโนโลยีและการวิเคราะห์" สำหรับรายงานชื่อ "{title}" จากบทความเหล่านี้:
{articles_text}
เขียนอย่างน้อย 4 ย่อหน้า พร้อมตารางเปรียบเทียบ"""),
            gen_section(f"""เขียน "ข้อเสนอแนะเชิงกลยุทธ์สำหรับ SYS" จากบทความเหล่านี้:
{articles_text}
เสนอแนะอย่างน้อย 6 ข้อ แต่ละข้อมีรายละเอียดการดำเนินการ"""),
            gen_section(f"""เขียน "บทสรุปและก้าวต่อไป" สำหรับรายงานชื่อ "{title}" จากบทความเหล่านี้:
{articles_text}
สรุปประเด็นสำคัญและเสนอขั้นตอนต่อไปที่เป็นรูปธรรม"""),
        )
        report = f"# {title}\n\n"
        report += f"## 1. บทสรุปผู้บริหาร\n\n{sections[0]}\n\n"
        report += f"## 2. ผลการวิจัยหลัก\n\n{sections[1]}\n\n"
        report += f"## 3. แนวโน้มเทคโนโลยีและการวิเคราะห์\n\n{sections[2]}\n\n"
        report += f"## 4. ข้อเสนอแนะเชิงกลยุทธ์สำหรับ SYS\n\n{sections[3]}\n\n"
        report += f"## 5. บทสรุปและก้าวต่อไป\n\n{sections[4]}\n"
    else:
        sections = await asyncio.gather(
            gen_section(f"""Write the "Executive Summary" section for a report titled "{title}".
Source articles:
{articles_text}
Write 4+ detailed paragraphs covering: overall significance, key themes, relevance to SYS, and what decision-makers need to know."""),
            gen_section(f"""Write the "Key Findings" section for a report titled "{title}".
Source articles:
{articles_text}
For each article write 2-3 paragraphs of detailed analysis. Then write 2 paragraphs synthesizing cross-cutting themes and patterns."""),
            gen_section(f"""Write the "Technology Trends & Analysis" section for a report titled "{title}".
Source articles:
{articles_text}
Write 4+ paragraphs analyzing technology trends and industry implications. Include a detailed comparison table with at least 5 rows comparing key aspects across articles."""),
            gen_section(f"""Write the "Strategic Recommendations for SYS" section for a report titled "{title}".
Source articles:
{articles_text}
Provide 6-8 specific, actionable recommendations. For each: bold title + 3 sentences explaining what to do, why it matters, and how to implement it."""),
            gen_section(f"""Write the "Conclusion & Next Steps" section for a report titled "{title}".
Source articles:
{articles_text}
Write 3 paragraphs: key takeaways, concrete next steps with timeline, and closing thoughts on strategic importance for SYS."""),
        )
        report = f"# {title}\n\n"
        report += f"## 1. Executive Summary\n\n{sections[0]}\n\n"
        report += f"## 2. Key Findings\n\n{sections[1]}\n\n"
        report += f"## 3. Technology Trends & Analysis\n\n{sections[2]}\n\n"
        report += f"## 4. Strategic Recommendations for SYS\n\n{sections[3]}\n\n"
        report += f"## 5. Conclusion & Next Steps\n\n{sections[4]}\n"

    return report


async def _llm_call_report(prompt: str, system: str) -> str:
    """LLM call specifically for report generation with higher token limit."""
    provider = settings.LLM_PROVIDER
    try:
        if "bedrock" in provider or provider == "bedrock":
            import boto3
            client = boto3.client(
                "bedrock-runtime",
                region_name=settings.AWS_REGION,
                aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
                aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            )
            body = json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": settings.REPORT_MAX_TOKENS,
                "system": system,
                "messages": [{"role": "user", "content": prompt}],
            })
            response = client.invoke_model(modelId=settings.BEDROCK_MODEL_ID, body=body)
            result = json.loads(response["body"].read())
            return result["content"][0]["text"]
        else:
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
            resp = await client.chat.completions.create(
                model=settings.OPENAI_MODEL,
                messages=[
                    {"role": "system", "content": system},
                    {"role": "user", "content": prompt},
                ],
                max_tokens=settings.REPORT_MAX_TOKENS,
                temperature=0.6,
                presence_penalty=0.2,
            )
            return resp.choices[0].message.content
    except Exception as e:
        return f"[Section generation error: {e}]"


async def generate_pptx(articles: list, title: str) -> str:
    """Generate a PowerPoint presentation from articles."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import tempfile

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"SYS Knowledge Hub | {len(articles)} Articles"

    # One slide per article
    for article in articles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = article.title[:80]
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"Category: {article.topic_category}"
        insights = (article.key_insights or [])[:5]
        for insight in insights:
            p = tf.add_paragraph()
            p.text = f"• {insight}"
            p.level = 1

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Topics Covered"
    content = slide.placeholders[1]
    tf = content.text_frame
    for a in articles:
        p = tf.add_paragraph()
        p.text = f"• {a.title[:60]}"

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name
